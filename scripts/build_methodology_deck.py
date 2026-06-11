#!/usr/bin/env python3
"""Populate Redrob idea submission PPTX from current team content."""

from __future__ import annotations

import argparse
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = PROJECT_ROOT / "challenge" / "Idea Submission Template _ Redrob.pptx"


def _set_text(shape, text: str, *, font_pt: int = 14) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_pt)


def build_deck(template: Path, output: Path) -> None:
    prs = Presentation(str(template))

    # Slide 1 — title fields (text boxes 1–3)
    s1 = prs.slides[0].shapes
    _set_text(s1[1], "Team Name : Sarva Automata", font_pt=20)
    _set_text(
        s1[2],
        "Problem Statement : Rank 100K candidates for Senior AI Engineer "
        "(search/retrieval/ranking); NDCG@10 focus; CPU-only; honeypot-aware.",
        font_pt=16,
    )
    _set_text(s1[3], "Team Leader Name : Arjun Mangarath", font_pt=18)

    slide_bodies = [
        # Slide 2
        "Hybrid recall (BM25 + BGE dense + career RRF) → LightGBM LambdaRank → "
        "BGE cross-encoder on top-700 with tuned fusion.\n\n"
        "Differentiation: IR + learned ranking + Redrob behavioral signals—not keywords. "
        "Explicit trap/honeypot and availability handling.",
        # Slide 3
        "JD: 5–8 YOE production ML, search/retrieval/ranking, embeddings, India hubs, "
        "product over research, recruiter reachability.\n\n"
        "Signals: silver tiers, LTR features, career evidence, GitHub, trap flags, "
        "behavioral twins (response, notice, recency). Cross-encoder for semantic fit—no runtime LLM.",
        # Slide 4
        "Retrieve ~2K (BM25 + FAISS + career RRF) → LTR + penalties → CE rerank top-700 → "
        "fusion + modifiers → top-100 CSV.\n\n"
        "Models: BM25, FAISS, LightGBM, BGE-reranker-base, Optuna-tuned fusion/modifiers.\n\n"
        "Combine: weighted CE/LTR/RRF × behavioral multiplier − trap/clone/research penalties.",
        # Slide 5
        "Reasoning: 5–6 templates from whitelisted fields (YOE, skills, location, availability, "
        "production vs research).\n\n"
        "No LLM at rank time → no hallucinated justifications.\n\n"
        "Traps: hard-exclude honeypots; clone cap top-30; availability down-rank; research penalty.",
        # Slide 6
        "Offline: preprocess → features/embeddings/BM25/FAISS → silver labels → train LTR → "
        "tune modifiers (~90 min).\n\n"
        "Online: rank.py loads artifacts → recall → LTR → CE → fusion/modifiers → "
        "team_sarva_automata.csv + audit.",
        # Slide 7
        "Offline: preprocess.py, train_ltr.py, tune_modifiers.py\n"
        "Online: rank.py + HF Docker sandbox\n"
        "Artifacts: parquet features, FAISS, BM25, ltr_model.lgb, fusion_params, modifier_params",
        # Slide 8
        "Audit: proxy NDCG@10 = 1.0; NDCG@30 = 0.9922; tier-5 top 10/30/100 = 9/26/70; "
        "0 honeypots in top-100; 0 low-availability in top-30.\n\n"
        "Constraints: precomputed indexes; rank ~5–8 min CPU, 16GB RAM, no network at inference.",
        # Slide 9
        "Python, Polars, LightGBM, sentence-transformers, FAISS, rank-bm25, FastAPI, Optuna. "
        "Gemini offline (JD + labels). Cursor for engineering assist.",
        # Slide 10
        "GitHub: github.com/ashokbugude/recruiter-candidate\n"
        "Sandbox: huggingface.co/spaces/ashokbugude/redrob-ranker "
        "(upload JSONL/JSON ≤100 → ranked CSV)\n"
        "Portal CSV: team_sarva_automata.csv\n"
        "Reproduce: python rank.py --candidates ./challenge/candidates.jsonl "
        "--out ./team_sarva_automata.csv",
    ]

    for slide_idx, body in enumerate(slide_bodies, start=1):
        shapes = prs.slides[slide_idx].shapes
        content_shape = shapes[-1]
        _set_text(content_shape, body, font_pt=13)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main() -> int:
    parser = argparse.ArgumentParser(description="Build methodology PPTX for portal upload.")
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument(
        "--out",
        type=Path,
        default=DEFAULT_TEMPLATE,
        help="Output .pptx (default: overwrite template in challenge/)",
    )
    parser.add_argument(
        "--export-pdf",
        type=Path,
        default=None,
        help="Optional PDF path via PowerPoint COM (Windows only)",
    )
    args = parser.parse_args()

    if args.out.resolve() != args.template.resolve():
        shutil.copy2(args.template, args.out)

    build_deck(args.template if args.out == args.template else args.out, args.out)
    print(f"Wrote deck: {args.out}")

    if args.export_pdf:
        try:
            import comtypes.client  # type: ignore

            pdf_path = args.export_pdf.resolve()
            pdf_path.parent.mkdir(parents=True, exist_ok=True)
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            deck = powerpoint.Presentations.Open(str(args.out.resolve()), WithWindow=False)
            deck.SaveAs(str(pdf_path), 32)
            deck.Close()
            powerpoint.Quit()
            print(f"Wrote PDF: {pdf_path}")
        except Exception as exc:
            print(f"PDF export skipped ({exc}). Export manually from PowerPoint: Save as PDF")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
